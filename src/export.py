"""Export utilities for CrystalBall visualizations.

This module handles the export of visualizations and analysis results to
various formats including PowerPoint presentations and PDF reports.
"""

import json
import os
from datetime import datetime
from typing import Dict, List, Optional, Union

import matplotlib.pyplot as plt
import pandas as pd
from docx import Document
from docx.shared import Inches as DocxInches
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches


def _normalize_filename_base(name: str, max_len: int = 120) -> str:
    """Normalize and shorten a base filename for exports.

    Steps:
    - Remove illegal filesystem characters
    - Split into tokens on underscores/whitespace
    - Apply conservative aliasing to common long tokens to shorten names
    - Lowercase all tokens for consistency
    - De-duplicate tokens while preserving order
    - Collapse repeated separators and enforce max length

    Example: 'consumerpriceinflationdetailedreferencetables_model_summary_rankings'
      -> 'cpi_detailed_ref_msum_rank'
    """
    import re
    if not name:
        return name

    # remove illegal filesystem characters and trim
    cleaned = re.sub(r'[<>:\"/\\|?*]+', '', name).strip('_ ').strip()

    # Split on underscores or whitespace
    raw_tokens = [t for t in re.split(r'[_\s]+', cleaned) if t]

    # Aliases to shorten common long tokens (all keys must be lowercase)
    alias = {
        # domains
        'consumerpriceinflationdetailedreferencetables': 'cpi_detailed_ref',
        'bank': 'bank',
        'rate': 'rate',
        'history': 'hist',
        'england': 'uk',
        'database': 'db',
        # pipeline/modeling
        'quartz': 'qz',
        'forecast': 'fcst',
        'model': 'mdl',
        'summary': 'sum',
        'modelsummary': 'msum',
        'model_summary': 'msum',
        'rank': 'rank',
        'rankings': 'rank',
        'rolling': 'roll',
        'corr': 'corr',
        'train': 'trn',
        'rmse': 'rmse',
        'affinity': 'aff',
        # misc
        'and': 'and',  # keep small connector to avoid accidental merges
        'data': 'data',
        'of': 'of',
    }

    # Map and lowercase
    mapped = []
    for t in raw_tokens:
        tl = t.lower()
        # allow aliases for combined tokens (e.g., 'model_summary')
        mapped_token = alias.get(tl, alias.get(tl.replace('-', ''), tl))
        mapped.append(mapped_token)

    # Post-process: collapse specific patterns
    # Convert mdl + sum -> msum
    combined = []
    i = 0
    while i < len(mapped):
        if i + 1 < len(mapped) and mapped[i] == 'mdl' and mapped[i + 1] == 'sum':
            combined.append('msum')
            i += 2
        else:
            combined.append(mapped[i])
            i += 1

    # Deduplicate while preserving order
    seen = set()
    collapsed: list[str] = []
    for t in combined:
        if t in seen:
            continue
        seen.add(t)
        collapsed.append(t)

    out = '_'.join(collapsed)
    out = re.sub(r'[_-]{2,}', '_', out).strip('_')

    if len(out) > max_len:
        out = out[:max_len].rstrip('_')
    return out


def create_forecast_slides(prs: Presentation,
                         models: Dict,
                         series_name: str,
                         rankings_df: pd.DataFrame,
                         image_paths: Dict[str, str]) -> Presentation:
    """
    Create slides for forecast visualizations.
    
    Args:
        prs: PowerPoint presentation object
        models: Dictionary of model results
        series_name: Name of the analyzed series
        rankings_df: DataFrame with model rankings
        image_paths: Dictionary mapping image types to file paths
        
    Returns:
        Updated presentation object
    """
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = f"Forecast Analysis: {series_name}"
    subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    
    # Rankings slide
    ranking_slide = prs.slides.add_slide(prs.slide_layouts[1])
    ranking_slide.shapes.title.text = "Model Rankings"
    
    # Add rankings table
    table_data = rankings_df.round(4)
    rows, cols = len(table_data) + 1, len(table_data.columns)
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(0.5) * (rows + 1)
    
    table = ranking_slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Headers
    for i, col in enumerate(table_data.columns):
        table.cell(0, i).text = col
        
    # Data
    for i, row in enumerate(table_data.itertuples(), start=1):
        for j, value in enumerate(row[1:]):
            table.cell(i, j).text = str(value)
            
    # Model performance slides (support flexible key names containing 'rankings')
    if any('rankings' in k.lower() for k in image_paths):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Model Performance Comparison"
        # Prefer an exact key if present, else take the first matching key
        img_key = 'rankings.png' if 'rankings.png' in image_paths else next((k for k in image_paths if 'rankings' in k.lower()), None)
        img_path = image_paths.get(img_key)
        left = Inches(1)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(img_path, left, top, height=Inches(5))
        
    # Forecast plots
    for model_name, path in image_paths.items():
        if 'forecast' in model_name.lower():
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"Forecast: {model_name}"
            left = Inches(1)
            top = Inches(1.5)
            pic = slide.shapes.add_picture(path, left, top, height=Inches(5))
            
    return prs

def export_to_powerpoint(models: Dict,
                        series_name: str,
                        rankings_df: pd.DataFrame,
                        image_paths: Dict[str, str],
                        output_path: str) -> str:
    """
    Export forecast analysis to PowerPoint presentation.
    
    Args:
        models: Dictionary of model results
        series_name: Name of the analyzed series
        rankings_df: DataFrame with model rankings
        image_paths: Dictionary mapping image types to file paths
        output_path: Path to save the PowerPoint file
        
    Returns:
        Path to the saved PowerPoint file
    """
    prs = Presentation()
    prs = create_forecast_slides(prs, models, series_name, rankings_df, image_paths)
    prs.save(output_path)
    return output_path

def export_to_pdf(models: Dict,
                 series_name: str,
                 rankings_df: pd.DataFrame,
                 image_paths: Dict[str, str],
                 output_path: str) -> str:
    """
    Export forecast analysis to PDF report.
    
    Args:
        models: Dictionary of model results
        series_name: Name of the analyzed series
        rankings_df: DataFrame with model rankings
        image_paths: Dictionary mapping image types to file paths
        output_path: Path to save the PDF file
        
    Returns:
        Path to the saved PDF file
    """
    class PDF(FPDF):
        def header(self):
            # Use core Helvetica to avoid fpdf2 Arial substitution deprecation
            self.set_font('helvetica', 'B', 12)
            try:
                from fpdf.enums import XPos, YPos
                self.cell(0, 10, 'CrystalBall Forecast Analysis', new_x=XPos.LMARGIN, new_y=YPos.NEXT, align='C')
            except Exception:
                self.cell(0, 10, 'CrystalBall Forecast Analysis', ln=True, align='C')
            self.ln(2)
            
        def footer(self):
            self.set_y(-15)
            self.set_font('helvetica', 'I', 8)
            try:
                from fpdf.enums import XPos, YPos
                self.cell(0, 10, f'Page {self.page_no()}', new_x=XPos.RIGHT, new_y=YPos.TOP, align='C')
            except Exception:
                self.cell(0, 10, f'Page {self.page_no()}', ln=0, align='C')
    
    pdf = PDF()
    pdf.add_page()
    
    # Title
    pdf.set_font('helvetica', 'B', 16)
    try:
        from fpdf.enums import XPos, YPos
        pdf.cell(0, 10, f"Forecast Analysis: {series_name}", new_x=XPos.LMARGIN, new_y=YPos.NEXT, align='C')
    except Exception:
        pdf.cell(0, 10, f"Forecast Analysis: {series_name}", ln=True, align='C')
    pdf.ln(10)
    
    # Rankings table
    pdf.set_font('helvetica', 'B', 12)
    try:
        from fpdf.enums import XPos, YPos
        pdf.cell(0, 10, "Model Rankings", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    except Exception:
        pdf.cell(0, 10, "Model Rankings", ln=True)
    pdf.ln(5)
    
    # Convert rankings to table
    table_data = rankings_df.round(4)
    col_width = pdf.w / len(table_data.columns)
    
    # Headers
    pdf.set_font('helvetica', 'B', 10)
    for col in table_data.columns:
        pdf.cell(col_width, 10, str(col), 1)
    pdf.ln()
    
    # Data
    pdf.set_font('helvetica', '', 10)
    for row in table_data.itertuples():
        for value in row[1:]:
            pdf.cell(col_width, 10, str(value), 1)
        pdf.ln()
        
    # Model performance plots
    for name, path in image_paths.items():
        pdf.add_page()
        pdf.set_font('helvetica', 'B', 12)
        if 'rankings' in name.lower():
            try:
                from fpdf.enums import XPos, YPos
                pdf.cell(0, 10, "Model Performance Comparison", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            except Exception:
                pdf.cell(0, 10, "Model Performance Comparison", ln=True)
        else:
            try:
                from fpdf.enums import XPos, YPos
                pdf.cell(0, 10, f"Forecast: {name}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            except Exception:
                pdf.cell(0, 10, f"Forecast: {name}", ln=True)
        
        # Embed image directly (PNG/JPG supported by FPDF)
        try:
            pdf.image(path, x=10, y=pdf.get_y(), w=190)
        except Exception:
            # Fallback: render via matplotlib if direct embedding fails
            plt.figure()
            img = plt.imread(path)
            plt.imshow(img)
            plt.axis('off')
            tmp_png = path + ".tmp.png"
            plt.savefig(tmp_png, format='png', bbox_inches='tight', pad_inches=0)
            plt.close()
            pdf.image(tmp_png, x=10, y=pdf.get_y(), w=190)
            if os.path.exists(tmp_png):
                os.remove(tmp_png)
        
        pdf.ln(10)
    
    pdf.output(output_path)
    return output_path

def export_to_word_docx(series_name: str,
                        rankings_df: pd.DataFrame,
                        image_paths: Dict[str, str],
                        output_path: str) -> str:
    """
    Export analysis to a Word DOCX document with rankings table and images.
    """
    doc = Document()
    doc.add_heading(f"Forecast Analysis: {series_name}", level=1)
    doc.add_paragraph(f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}")

    # Rankings table
    doc.add_heading("Model Rankings", level=2)
    tbl = doc.add_table(rows=1, cols=len(rankings_df.columns))
    hdr_cells = tbl.rows[0].cells
    for i, col in enumerate(rankings_df.columns):
        hdr_cells[i].text = str(col)
    for _, row in rankings_df.iterrows():
        cells = tbl.add_row().cells
        for i, col in enumerate(rankings_df.columns):
            cells[i].text = str(row[col])

    # Images
    doc.add_page_break()
    if 'rankings' in ''.join(image_paths.keys()).lower():
        for name, path in image_paths.items():
            if 'rankings' in name.lower() and os.path.exists(path):
                doc.add_heading("Model Performance Comparison", level=2)
                doc.add_picture(path, width=DocxInches(6.5))
                break

    # Forecast plots
    for name, path in image_paths.items():
        if 'forecast' in name.lower() and os.path.exists(path):
            doc.add_heading(f"Forecast: {name}", level=3)
            doc.add_picture(path, width=DocxInches(6.5))

    doc.save(output_path)
    return output_path


def export_to_excel(series_name: str,
                    rankings_df: pd.DataFrame,
                    image_paths: Dict[str, str],
                    output_path: str) -> str:
    """
    Export analysis to an Excel workbook with rankings sheet and embedded chart/images.
    """
    with pd.ExcelWriter(output_path, engine='xlsxwriter') as writer:
        rankings_df.to_excel(writer, index=False, sheet_name='Rankings')
        wb = writer.book
        ws = writer.sheets['Rankings']
        # Create a bar chart for weighted_score if present, else mean_rank if present
        chart = wb.add_chart({'type': 'bar'})
        if 'weighted_score' in rankings_df.columns:
            cat_col = rankings_df.columns.get_loc('model')
            val_col = rankings_df.columns.get_loc('weighted_score')
            chart.add_series({
                'name': 'Weighted Score',
                'categories': ['Rankings', 1, cat_col, len(rankings_df), cat_col],
                'values':     ['Rankings', 1, val_col, len(rankings_df), val_col],
            })
            chart.set_title({'name': f'Weighted Score - {series_name}'})
            chart.set_x_axis({'name': 'Model'})
            chart.set_y_axis({'name': 'Weighted Score'})
        elif 'mean_rank' in rankings_df.columns:
            cat_col = rankings_df.columns.get_loc('model')
            val_col = rankings_df.columns.get_loc('mean_rank')
            chart.add_series({
                'name': 'Mean Rank',
                'categories': ['Rankings', 1, cat_col, len(rankings_df), cat_col],
                'values':     ['Rankings', 1, val_col, len(rankings_df), val_col],
            })
            chart.set_title({'name': f'Mean Rank - {series_name}'})
            chart.set_x_axis({'name': 'Model'})
            chart.set_y_axis({'name': 'Mean Rank'})
        chart.set_legend({'position': 'bottom'})
        ws.insert_chart('G2', chart, {'x_scale': 1.2, 'y_scale': 1.2})

        # Optional: add images to a second sheet
        img_sheet = wb.add_worksheet('Images')
        r, c = 1, 1
        for name, path in image_paths.items():
            if os.path.exists(path) and path.lower().endswith('.png'):
                img_sheet.write(r, c, name)
                img_sheet.insert_image(r + 1, c, path, {'x_scale': 0.7, 'y_scale': 0.7})
                r += 20
    return output_path


def batch_export_reports(results_dir: str,
                        output_dir: str,
                        formats: List[str] = ['pptx', 'pdf', 'docx', 'xlsx']) -> Dict[str, List[str]]:
    """
    Batch export all analysis results to PowerPoint and/or PDF.
    
    Args:
        results_dir: Directory containing analysis results
        output_dir: Directory to save exported files
        formats: List of formats to export ('pptx' and/or 'pdf')
        
    Returns:
        Dictionary mapping formats to lists of exported file paths
    """
    os.makedirs(output_dir, exist_ok=True)
    exported_files = {fmt: [] for fmt in formats}
    # Track collisions on normalized base names to avoid overwriting
    collisions: Dict[str, int] = {}

    # Determine base names from rankings CSVs
    ranking_files = [f for f in os.listdir(results_dir) if f.endswith('_rankings.csv')]
    for rf in ranking_files:
        base_name = rf.replace('_rankings.csv', '')
        norm_base = _normalize_filename_base(base_name)
        # Disambiguate if this normalized base was already used
        if norm_base in collisions:
            collisions[norm_base] += 1
            final_base = f"{norm_base}_{collisions[norm_base]}"
        else:
            collisions[norm_base] = 1
            final_base = norm_base

        # Load rankings
        rankings_path = os.path.join(results_dir, rf)
        try:
            rankings_df = pd.read_csv(rankings_path)
        except Exception:
            continue

        # Collect image paths from visuals directory
        vis_dir = os.path.join(results_dir, 'visualizations')
        image_paths: Dict[str, str] = {}
        if os.path.isdir(vis_dir):
            for img in os.listdir(vis_dir):
                if img.startswith(base_name) and img.lower().endswith('.png'):
                    key = img.replace(f"{base_name}_", "").replace(".png", "")
                    image_paths[key] = os.path.join(vis_dir, img)

        # Export to each format
        if 'pptx' in formats:
            pptx_path = os.path.join(output_dir, f"{final_base}_analysis.pptx")
            try:
                export_to_powerpoint(
                    {},
                    base_name,
                    rankings_df,
                    image_paths,
                    pptx_path
                )
                exported_files['pptx'].append(pptx_path)
            except Exception as e:
                print(f"Failed to export PowerPoint for {base_name}: {str(e)}")
                
        if 'pdf' in formats:
            pdf_path = os.path.join(output_dir, f"{final_base}_analysis.pdf")
            try:
                export_to_pdf(
                    {},
                    base_name,
                    rankings_df,
                    image_paths,
                    pdf_path
                )
                exported_files['pdf'].append(pdf_path)
            except Exception as e:
                print(f"Failed to export PDF for {base_name}: {str(e)}")

        if 'docx' in formats:
            docx_path = os.path.join(output_dir, f"{final_base}_analysis.docx")
            try:
                export_to_word_docx(
                    base_name,
                    rankings_df,
                    image_paths,
                    docx_path,
                )
                exported_files['docx'].append(docx_path)
            except Exception as e:
                print(f"Failed to export DOCX for {base_name}: {str(e)}")

        if 'xlsx' in formats:
            xlsx_path = os.path.join(output_dir, f"{final_base}_analysis.xlsx")
            try:
                export_to_excel(
                    base_name,
                    rankings_df,
                    image_paths,
                    xlsx_path,
                )
                exported_files['xlsx'].append(xlsx_path)
            except Exception as e:
                print(f"Failed to export XLSX for {base_name}: {str(e)}")
                    
    return exported_files


def clear_duplicate_exports(output_dir: str, formats: Optional[List[str]] = None) -> Dict[str, List[str]]:
    """Remove duplicate export files, keeping only the most recently modified in each group.

    Duplicates are defined as files whose names differ only by a trailing numeric suffix
    immediately before the "_analysis" marker (e.g., "..._rankings_2_analysis.pdf").
    We group such files by their base name with that numeric suffix removed, and keep only
    the most recent file per group. Files that do not have a corresponding base without the
    numeric suffix are left untouched to avoid deleting legitimate numbered series like
    "series-160925_1_analysis".

    Args:
        output_dir: Directory containing exported files
        formats: Optional list of formats (extensions) to consider (e.g., ['pdf','pptx']).

    Returns:
        Dict with keys 'kept' and 'removed' listing affected file paths.
    """
    import re
    kept: List[str] = []
    removed: List[str] = []

    if not os.path.isdir(output_dir):
        return {'kept': kept, 'removed': removed}

    # Default to all known export extensions
    exts = formats or ['pdf', 'pptx', 'docx', 'xlsx']
    exts = [e.lower().lstrip('.') for e in exts]

    # Gather files by extension and build set of stems (without extension)
    all_files: List[str] = []
    stems: set[str] = set()
    for fname in os.listdir(output_dir):
        ext = os.path.splitext(fname)[1].lower().lstrip('.')
        if ext in exts:
            fpath = os.path.join(output_dir, fname)
            if os.path.isfile(fpath):
                all_files.append(fpath)
                stems.add(os.path.splitext(fname)[0])

    # Build groups keyed by canonical stem
    groups: Dict[str, List[str]] = {}
    for fpath in all_files:
        stem = os.path.splitext(os.path.basename(fpath))[0]
        # Try to strip a trailing _<digits> only if a base without it also exists
        stripped = re.sub(r'_(\d+)(?=_analysis$)', '', stem)
        if stripped != stem and stripped in stems:
            key = stripped
        else:
            key = stem
        groups.setdefault(key, []).append(fpath)

    # For each group, if more than one file, keep the most recent and remove the others
    for key, files in groups.items():
        if len(files) <= 1:
            kept.extend(files)
            continue
        files_sorted = sorted(files, key=lambda p: os.path.getmtime(p), reverse=True)
        kept.append(files_sorted[0])
        for old in files_sorted[1:]:
            try:
                os.remove(old)
                removed.append(old)
            except Exception:
                # If removal fails, keep the file to avoid data loss
                kept.append(old)

    return {'kept': kept, 'removed': removed}


def generate_exports_manifest(output_dir: str, manifest_path: Optional[str] = None,
                              formats: Optional[List[str]] = None) -> Dict[str, Union[str, List[Dict[str, Union[str, int]]]]]:
    """Create a manifest JSON for current export artifacts with timestamps and groups.

    The grouping logic mirrors duplicate detection in clear_duplicate_exports: files with
    trailing _<digits> before _analysis.(ext) are grouped with their base if present.

    Args:
        output_dir: Directory containing exported files.
        manifest_path: If provided, write JSON manifest to this path.
        formats: Optional subset of extensions to include.

    Returns:
        Dict containing 'generated_at', 'directory', 'file_count', 'group_count',
        'files' list, and 'groups' summary.
    """
    import re
    from datetime import datetime as _dt

    result: Dict[str, Union[str, int, List[Dict[str, Union[str, int]]]]] = {}
    if not os.path.isdir(output_dir):
        result = {
            'generated_at': _dt.utcnow().isoformat() + 'Z',
            'directory': output_dir,
            'file_count': 0,
            'group_count': 0,
            'files': [],
            'groups': {},
        }
        if manifest_path:
            with open(manifest_path, 'w', encoding='utf-8') as f:
                json.dump(result, f, indent=2)
        return result

    exts = [e.lower().lstrip('.') for e in (formats or ['pdf', 'pptx', 'docx', 'xlsx'])]

    # Collect files
    entries: List[Dict[str, Union[str, int]]] = []
    stems: set[str] = set()
    for fname in os.listdir(output_dir):
        ext = os.path.splitext(fname)[1].lower().lstrip('.')
        if ext not in exts:
            continue
        fpath = os.path.join(output_dir, fname)
        if not os.path.isfile(fpath):
            continue
        stat = os.stat(fpath)
        mtime = stat.st_mtime
        iso = _dt.fromtimestamp(mtime).isoformat()
        size = stat.st_size
        stem = os.path.splitext(fname)[0]
        stems.add(stem)
        entries.append({
            'name': fname,
            'path': fpath,
            'ext': ext,
            'size_bytes': size,
            'modified': iso,
        })

    # Build groups based on same logic as duplicate clearance
    groups: Dict[str, List[Dict[str, Union[str, int]]]] = {}
    for item in entries:
        stem = os.path.splitext(item['name'])[0]  # type: ignore[index]
        stripped = re.sub(r'_(\d+)(?=_analysis$)', '', stem)
        key = stripped if (stripped != stem and stripped in stems) else stem
        groups.setdefault(key, []).append(item)

    # Sort entries within groups by modified desc
    for key, items in groups.items():
        items.sort(key=lambda x: x['modified'], reverse=True)  # type: ignore[index]

    manifest = {
        'generated_at': _dt.utcnow().isoformat() + 'Z',
        'directory': output_dir,
        'file_count': len(entries),
        'group_count': len(groups),
        'files': sorted(entries, key=lambda x: x['name']),  # type: ignore[index]
        'groups': {k: [{'name': i['name'], 'modified': i['modified'], 'size_bytes': i['size_bytes'], 'ext': i['ext']} for i in v] for k, v in sorted(groups.items())},
    }

    if manifest_path:
        with open(manifest_path, 'w', encoding='utf-8') as f:
            json.dump(manifest, f, indent=2)
    return manifest